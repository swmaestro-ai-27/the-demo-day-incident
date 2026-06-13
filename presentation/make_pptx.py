import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

def run():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Color Constants
    COLOR_BG = RGBColor(2, 3, 5)
    COLOR_TEXT = RGBColor(244, 244, 245)
    COLOR_MUTED = RGBColor(156, 163, 175)
    COLOR_DIM = RGBColor(82, 88, 102)
    COLOR_DANGER = RGBColor(255, 51, 71)
    COLOR_AMBER = RGBColor(246, 193, 91)
    COLOR_CARD_BG = RGBColor(8, 9, 11)
    COLOR_LINE = RGBColor(37, 40, 50)

    # Helpers
    def set_slide_background(slide, color):
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = color

    def add_card(slide, left, top, width, height, active=False):
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
        shape.fill.solid()
        shape.fill.fore_color.rgb = COLOR_CARD_BG
        shape.line.color.rgb = COLOR_DANGER if active else COLOR_LINE
        shape.line.width = Pt(1.5) if active else Pt(1.0)
        return shape

    def add_text_box(slide, left, top, width, height, text, font_size, bold=False, color=COLOR_TEXT, align=PP_ALIGN.LEFT):
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        tf.word_wrap = True
        tf.margin_left = Inches(0)
        tf.margin_top = Inches(0)
        tf.margin_right = Inches(0)
        tf.margin_bottom = Inches(0)
        p = tf.paragraphs[0]
        p.text = text
        p.alignment = align
        p.font.size = Pt(font_size)
        p.font.bold = bold
        p.font.color.rgb = color
        p.font.name = "Malgun Gothic"
        return txBox

    def add_formatted_text(slide, left, top, width, height, title, body, title_size=15, body_size=11, title_color=COLOR_TEXT, body_color=COLOR_MUTED):
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        tf.word_wrap = True
        tf.margin_left = Inches(0)
        tf.margin_top = Inches(0)
        tf.margin_right = Inches(0)
        tf.margin_bottom = Inches(0)
        
        # Title paragraph
        p1 = tf.paragraphs[0]
        p1.text = title
        p1.font.size = Pt(title_size)
        p1.font.bold = True
        p1.font.color.rgb = title_color
        p1.font.name = "Malgun Gothic"
        p1.space_after = Pt(6)
        
        # Body paragraph
        p2 = tf.add_paragraph()
        p2.text = body
        p2.font.size = Pt(body_size)
        p2.font.color.rgb = body_color
        p2.font.name = "Malgun Gothic"
        p2.line_spacing = 1.3
        
        return txBox

    def create_slide(prs, title, eyebrow, progress_text):
        slide_layout = prs.slide_layouts[6]  # blank layout
        slide = prs.slides.add_slide(slide_layout)
        set_slide_background(slide, COLOR_BG)
        
        # Eyebrow
        add_text_box(slide, Inches(1.0), Inches(0.5), Inches(11.33), Inches(0.4), eyebrow, 11, bold=True, color=COLOR_DIM)
        
        # Title
        add_text_box(slide, Inches(1.0), Inches(0.9), Inches(11.33), Inches(0.8), title, 32, bold=True, color=COLOR_TEXT)
        
        # Progress footer
        add_text_box(slide, Inches(1.0), Inches(6.8), Inches(11.33), Inches(0.4), progress_text, 10, bold=True, color=COLOR_DIM, align=PP_ALIGN.RIGHT)
        
        return slide

    def add_flow_node(slide, left, top, width, height, text, active=False):
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
        shape.fill.solid()
        if active:
            shape.fill.fore_color.rgb = RGBColor(25, 12, 14)
            shape.line.color.rgb = COLOR_DANGER
            shape.line.width = Pt(1.5)
            color = COLOR_DANGER
        else:
            shape.fill.fore_color.rgb = COLOR_CARD_BG
            shape.line.color.rgb = COLOR_LINE
            shape.line.width = Pt(1.0)
            color = COLOR_TEXT
        
        tf = shape.text_frame
        tf.word_wrap = True
        tf.margin_left = Inches(0.1)
        tf.margin_right = Inches(0.1)
        tf.margin_top = Inches(0.05)
        tf.margin_bottom = Inches(0.05)
        p = tf.paragraphs[0]
        p.text = text
        p.alignment = PP_ALIGN.CENTER
        p.font.size = Pt(11)
        p.font.bold = True
        p.font.color.rgb = color
        p.font.name = "Malgun Gothic"

    def add_browser_mockup(slide, left, top, width, image_path):
        height = int(width / 1.6)
        header_height = Inches(0.22)
        
        # 1. Outer Container Shape (rounded rectangle)
        outer = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height + header_height)
        outer.fill.solid()
        outer.fill.fore_color.rgb = COLOR_CARD_BG
        outer.line.color.rgb = COLOR_LINE
        outer.line.width = Pt(1.0)
        
        # 2. Add the Image inside the container (below the header)
        slide.shapes.add_picture(image_path, left + Inches(0.01), top + header_height, width - Inches(0.02), height - Inches(0.01))
        
        # 3. Add three dots for browser control (red, yellow, green)
        dot_size = Inches(0.06)
        dot_y = top + Inches(0.08)
        dot_colors = [RGBColor(255, 95, 86), RGBColor(255, 189, 46), RGBColor(39, 201, 63)]
        for idx, color in enumerate(dot_colors):
            dot_x = left + Inches(0.12) + (idx * Inches(0.1))
            dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, dot_x, dot_y, dot_size, dot_size)
            dot.fill.solid()
            dot.fill.fore_color.rgb = color
            dot.line.fill.background() # no border

    # ==========================================
    # Slide 1: Introduction
    # ==========================================
    slide1 = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide1, COLOR_BG)
    add_text_box(slide1, Inches(1.0), Inches(0.5), Inches(11.33), Inches(0.4), "CASE 01", 11, bold=True, color=COLOR_DIM)
    add_text_box(slide1, Inches(1.0), Inches(1.8), Inches(7.0), Inches(1.5), "The Demo Day Incident", 54, bold=True, color=COLOR_TEXT)
    add_text_box(slide1, Inches(1.0), Inches(3.4), Inches(6.8), Inches(1.0), "서로 다른 권한과 페르소나를 가진 에이전트들과 상호작용하며 사건을 해결하는 추리 게임", 18, color=COLOR_MUTED)
    
    # Quote Card
    add_card(slide1, Inches(1.0), Inches(4.7), Inches(6.5), Inches(1.3))
    add_text_box(slide1, Inches(1.3), Inches(4.9), Inches(5.9), Inches(0.9), "안녕하십니까?\n저희 조의 발표를 시작하겠습니다.\n\"NPC가 단서를 주는 사람이 아니라, 자기 나름대로 범인을 찍는 사람이라면?\"", 13, bold=True, color=COLOR_TEXT)
    
    # Home Screen Mockup
    if os.path.exists("presentation/assets/deck-home.png"):
        add_browser_mockup(slide1, Inches(8.1), Inches(2.2), Inches(4.2), "presentation/assets/deck-home.png")
    
    add_text_box(slide1, Inches(1.0), Inches(6.8), Inches(11.33), Inches(0.4), "01 / 06", 10, bold=True, color=COLOR_DIM, align=PP_ALIGN.RIGHT)

    # ==========================================
    # Slide 2: 1. 서비스 선정 배경 및 핵심 가치
    # ==========================================
    slide2 = create_slide(prs, "1. 서비스 선정 배경 및 핵심 가치", "PART 01", "02 / 06")
    
    # Card 1: 기존 게임의 아쉬움
    add_card(slide2, Inches(1.0), Inches(1.9), Inches(5.0), Inches(3.8))
    add_formatted_text(
        slide2, Inches(1.3), Inches(2.2), Inches(4.4), Inches(3.2),
        "기존 게임의 아쉬움",
        "• 수동적인 스토리라인\n  정해진 선택지 중 하나를 골라가며 추리를 진행해야 하기에 플레이어가 스스로 생각하기보다 짜인 스토리를 수동적으로 따라가는 한계가 있습니다.\n\n• 단순 정보 전달자\n  최근 AI 에이전트를 도입한 추리 게임 역시 에이전트들이 단순히 단서를 보관하고 일방적으로 전달하는 단순 도구 역할에 머무르고 있습니다."
    )
    
    # Card 2: 우리 서비스의 핵심 가치
    add_card(slide2, Inches(6.4), Inches(1.9), Inches(5.0), Inches(3.8))
    add_formatted_text(
        slide2, Inches(6.7), Inches(2.2), Inches(4.4), Inches(3.2),
        "우리 서비스의 핵심 가치",
        "• 능동적인 불완전 추리자\n  각 에이전트들이 자신에게 할당된 고유한 시스템 권한과 페르소나만을 바탕으로 사건을 해석하며, 서로를 능동적으로 의심하고 불완전한 추리를 전개합니다.\n\n• 진짜 근거를 가려내는 재미\n  이들의 그럴듯한 주장 속 모순과 객관적 로그 단서를 대조하여 플레이어가 스스로 진짜 증거를 가려내는 추리 본연의 재미를 선사합니다."
    )
    
    # Bottom Quote
    add_text_box(slide2, Inches(1.0), Inches(6.0), Inches(11.33), Inches(0.5), "\"AI가 범인을 맞히는 게임이 아니라, 플레이어가 AI의 추리를 의심하는 게임\"", 17, bold=True, color=COLOR_TEXT, align=PP_ALIGN.CENTER)

    # ==========================================
    # Slide 3: 2. 서비스 핵심 기능
    # ==========================================
    slide3 = create_slide(prs, "2. 서비스 핵심 기능", "PART 02", "03 / 06")
    
    # Steps
    steps = [
        ("사건 접속 & 단서 확보", "데모데이 전날 밤 연수생 의문사 현장 접속 및 서버 로그, 삭제 기록 등의 단서 잠금 해제"),
        ("실시간 에이전트 심문", "해금한 단서를 기반으로 용의자 에이전트들과 자유로운 채팅 대화로 진술 및 모순 포착"),
        ("추리 제출 및 판정", "단서와 주장을 비교해 범인(오케스트레이터 ARIA), 동기, 핵심 증거를 최종 제출하여 성공 판정 피드백 수신")
    ]
    for i, (title, desc) in enumerate(steps):
        top_pos = Inches(1.9 + (i * 1.5))
        add_card(slide3, Inches(1.0), top_pos, Inches(6.5), Inches(1.2))
        add_formatted_text(slide3, Inches(1.3), top_pos + Inches(0.2), Inches(5.9), Inches(0.8), f"{i+1}. {title}", desc, title_size=14, body_size=10.5)

    if os.path.exists("presentation/assets/deck-main.png"):
        add_browser_mockup(slide3, Inches(8.1), Inches(2.2), Inches(4.2), "presentation/assets/deck-main.png")

    # ==========================================
    # Slide 4: 3. Agent Workflow 기획 및 구성
    # ==========================================
    slide4 = create_slide(prs, "3. Agent Workflow 기획 및 구성", "PART 03", "04 / 06")
    
    # Left explanations
    add_card(slide4, Inches(1.0), Inches(1.9), Inches(6.0), Inches(2.1))
    add_formatted_text(
        slide4, Inches(1.3), Inches(2.1), Inches(5.4), Inches(1.7),
        "지식 필터링 (Knowledge Filtering)",
        "용의자 에이전트 지식 중 플레이어가 게임 진행을 통해 실제로 해금한 단서(context_clues)만 LLM 프롬프트에 동적으로 바인딩하여 주입합니다. 이를 통해 용의자가 자신의 정보 권한을 넘는 발언을 하지 못하도록 설계했습니다."
    )
    
    add_card(slide4, Inches(1.0), Inches(4.3), Inches(6.0), Inches(2.1))
    add_formatted_text(
        slide4, Inches(1.3), Inches(4.5), Inches(5.4), Inches(1.7),
        "실제 구현된 스포일러 가드 (Spoiler Guard)",
        "용의자가 대답하는 과정에서 미해금 단서(locked_clues)의 이름이나 핵심 키워드를 스포일러하는지 검사합니다. 유출 감지 시 안전한 우회 답변(\"지금 확인된 정보만으로는...\")으로 즉시 강제 대체하는 필터가 실제로 연동되어 동작합니다. (backend/agents/guard.py)"
    )

    # Right Flowchart
    add_flow_node(slide4, Inches(7.8), Inches(1.9), Inches(4.5), Inches(0.5), "사용자 입력 & 캐릭터 ID 수신")
    add_text_box(slide4, Inches(7.8), Inches(2.4), Inches(4.5), Inches(0.25), "↓", 11, bold=True, color=COLOR_DIM, align=PP_ALIGN.CENTER)
    
    add_flow_node(slide4, Inches(7.8), Inches(2.65), Inches(4.5), Inches(0.5), "해금 정보 대조 및 지식 필터링")
    add_text_box(slide4, Inches(7.8), Inches(3.15), Inches(4.5), Inches(0.25), "↓", 11, bold=True, color=COLOR_DIM, align=PP_ALIGN.CENTER)
    
    add_flow_node(slide4, Inches(7.8), Inches(3.4), Inches(4.5), Inches(0.5), "페르소나 기반 LLM 답변 생성")
    add_text_box(slide4, Inches(7.8), Inches(3.9), Inches(4.5), Inches(0.25), "↓", 11, bold=True, color=COLOR_DIM, align=PP_ALIGN.CENTER)
    
    add_flow_node(slide4, Inches(7.8), Inches(4.15), Inches(4.5), Inches(0.5), "스포일러 가드 (Spoiler Guard) 검사", active=True)
    
    # Branches
    add_text_box(slide4, Inches(7.8), Inches(4.65), Inches(2.2), Inches(0.25), "↙ [유출 감지]", 9, bold=True, color=COLOR_DANGER, align=PP_ALIGN.CENTER)
    add_text_box(slide4, Inches(10.1), Inches(4.65), Inches(2.2), Inches(0.25), "[안전함] ↘", 9, bold=True, color=COLOR_MUTED, align=PP_ALIGN.CENTER)
    
    add_flow_node(slide4, Inches(7.8), Inches(4.9), Inches(2.2), Inches(0.5), "우회형 답변 대체")
    add_flow_node(slide4, Inches(10.1), Inches(4.9), Inches(2.2), Inches(0.5), "원문 진술 반환")
    
    # Merge arrow
    add_text_box(slide4, Inches(7.8), Inches(5.4), Inches(4.5), Inches(0.25), "↓", 11, bold=True, color=COLOR_DIM, align=PP_ALIGN.CENTER)
    add_flow_node(slide4, Inches(7.8), Inches(5.65), Inches(4.5), Inches(0.5), "최종 응답 반환 및 대화 DB 저장")

    # ==========================================
    # Slide 5: 4. 발전 계획
    # ==========================================
    slide5 = create_slide(prs, "4. 발전 계획", "PART 04", "05 / 06")
    
    # Card 1: 자율적인 상호 소통 구조
    add_card(slide5, Inches(1.0), Inches(1.9), Inches(3.3), Inches(4.5))
    add_formatted_text(
        slide5, Inches(1.2), Inches(2.2), Inches(2.9), Inches(3.9),
        "자율적인 상호 소통 구조",
        "• 에이전트 간 소통\n  플레이어의 개입 없이도 에이전트들끼리 정보를 교환하고 서로 대화하여 상황과 추리가 실시간으로 역동적으로 변화하는 한층 더 복잡한 게임 환경을 구축합니다.\n\n• 실시간 정보 전이\n  한 에이전트의 비밀이 누설되면 다른 에이전트가 그 정보를 인지해 대화 내용에 실시간 반영되는 시스템을 구현합니다."
    )
    
    # Card 2: 시나리오 자동화 (Agentic Workflow)
    add_card(slide5, Inches(4.6), Inches(1.9), Inches(3.3), Inches(4.5))
    add_formatted_text(
        slide5, Inches(4.8), Inches(2.2), Inches(2.9), Inches(3.9),
        "시나리오 자동화",
        "• Agentic Workflow\n  이번 사건의 완성도 높은 시나리오 제작 프로세스를 템플릿화하여 시나리오 생성 에이전트를 구축합니다.\n\n• 완전 자동 시나리오 시스템\n  사건의 전말, 단서 배치, 용의자 정보 권한 및 상반된 알리바이까지 자동으로 완벽하게 빌드 및 배포되는 게임 시나리오 파이프라인을 기획 및 설계합니다."
    )

    if os.path.exists("presentation/assets/deck-character-chat.png"):
        add_browser_mockup(slide5, Inches(8.1), Inches(2.2), Inches(4.2), "presentation/assets/deck-character-chat.png")

    # ==========================================
    # Slide 6: 5. Agent Workflow 시연
    # ==========================================
    slide6 = create_slide(prs, "5. Agent Workflow 시연", "PART 05", "06 / 06")
    
    # Steps
    demo_steps = [
        ("현장 접속 및 조사", "오케스트레이터의 이상 개입 정황 발견 후 용의자 에이전트 심문"),
        ("에이전트 모순 진술", "각자의 시스템 권한 범위에서 서로 다른 주장을 하는 모순 포착"),
        ("최종 범인 제출", "시스템 오케스트레이터 ARIA를 지목하고 성공 판정 피드백 확인")
    ]
    for i, (title, desc) in enumerate(demo_steps):
        top_pos = Inches(1.9 + (i * 1.2))
        add_card(slide6, Inches(1.0), top_pos, Inches(6.0), Inches(1.0))
        add_formatted_text(slide6, Inches(1.2), top_pos + Inches(0.15), Inches(5.6), Inches(0.7), title, desc, title_size=13, body_size=10)

    # QR Code Panel
    add_card(slide6, Inches(1.0), Inches(5.5), Inches(6.0), Inches(1.1))
    if os.path.exists("presentation/assets/play-qr.png"):
        slide6.shapes.add_picture("presentation/assets/play-qr.png", Inches(1.1), Inches(5.55), Inches(1.0), Inches(1.0))
    add_formatted_text(
        slide6, Inches(2.3), Inches(5.65), Inches(4.5), Inches(0.8),
        "지금 바로 플레이",
        "frontend-plum-seven-77.vercel.app",
        title_size=14, body_size=11, title_color=COLOR_TEXT, body_color=COLOR_MUTED
    )

    # VIDEO EMBEDDING!
    video_path = "presentation/agent-workflow-demo.mp4"
    if os.path.exists(video_path):
        print(f"Embedding video: {video_path}")
        # Add video: left, top, width, height. Mime type is video/mp4.
        # Ratio of widescreen video is 16:9, so 4.8 width -> 2.7 height
        slide6.shapes.add_movie(
            video_path,
            Inches(7.6),
            Inches(2.1),
            Inches(4.6),
            Inches(2.6),
            mime_type='video/mp4'
        )
        
        # Add frame around video for premium aesthetic
        frame = slide6.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(7.6), Inches(2.1), Inches(4.6), Inches(2.6))
        frame.fill.background()
        frame.line.color.rgb = COLOR_LINE
        frame.line.width = Pt(1.5)
        
        # Video title text
        add_text_box(slide6, Inches(7.6), Inches(4.8), Inches(4.6), Inches(0.4), "* 데모 시연 영상 (agent-workflow-demo.mp4)", 11, color=COLOR_DIM, align=PP_ALIGN.CENTER)
    else:
        print("Video file not found!")
        # Fallback card
        add_card(slide6, Inches(7.6), Inches(2.1), Inches(4.6), Inches(3.0))
        add_text_box(slide6, Inches(7.9), Inches(3.2), Inches(4.0), Inches(1.0), "시연 비디오 준비 중\n(agent-workflow-demo.mp4)", 14, bold=True, color=COLOR_MUTED, align=PP_ALIGN.CENTER)

    # Save
    prs.save("presentation/demo-day-incident-deck.pptx")
    print("Successfully generated: presentation/demo-day-incident-deck.pptx")

if __name__ == "__main__":
    run()
